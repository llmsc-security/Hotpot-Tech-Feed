"""Build docs/slides.pptx from a small inline outline. Run:  python3 docs/build_slides.py"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "slides.pptx"

DARK = RGBColor(0x1F, 0x1F, 0x1F)
RED = RGBColor(0xB9, 0x1C, 0x1C)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CODE_BG = RGBColor(0x0F, 0x17, 0x2A)
CODE_FG = RGBColor(0xE2, 0xE8, 0xF0)
PANEL_BG = RGBColor(0xF8, 0xFA, 0xFC)
PANEL_BORDER = RGBColor(0xE2, 0xE8, 0xF0)


def add_text(slide, left, top, width, height, text, *, size=18,
             color=TEXT, bold=False, italic=False, align="left"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}[align]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def add_title(slide, text, *, size=34, color=DARK):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(11.8), Inches(0.9),
             text, size=size, bold=True, color=color)
    # Amber underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.2),
                                  Inches(0.9), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()


def add_panel(slide, left, top, width, height, *, bg=PANEL_BG, border=PANEL_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_code_block(slide, left, top, width, height, code):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = CODE_BG
    panel.line.fill.background()
    panel.shadow.inherit = False

    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = 1
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(14)
        run.font.name = "Consolas"
        run.font.color.rgb = CODE_FG


def add_bullets(slide, left, top, width, height, items, *, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "• " + text
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = color
    return box


def page_number(slide, n, total):
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.3),
             f"{n} / {total}", size=10, color=MUTED, align="right")


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Big amber bar on the left
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0), Inches(0.4), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()

    add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
             "🌶️ Hotpot Tech Feed", size=60, bold=True, color=DARK)
    add_text(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.7),
             "A daily CS digest, driven by an LLM agent",
             size=26, color=RED)
    add_text(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.6),
             "Self-hosted feed reader that ingests papers, blogs, and lab "
             "announcements,\nclassifies them with Qwen3.5, and answers "
             "natural-language search queries.",
             size=18, color=TEXT)
    add_text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.4),
             "Stack  ·  FastAPI · Postgres · Redis · Qdrant · React · nginx · Qwen3.5",
             size=14, color=MUTED)
    add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
             "Repo  ·  github.com/llmsc-security/Hotpot-Tech-Feed",
             size=14, color=MUTED)


def slide_architecture(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "1 · System architecture")

    # Left column: bullets
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(6.8), Inches(5.2), [
        "nginx (port 50002) is the only host-facing service; it serves the React SPA and reverse-proxies /api/* to the backend over an internal docker network.",
        "FastAPI runs the REST API and the hotpot CLI (ingest-now, enrich-all, send-test-digest, …).",
        "Postgres 16 holds items, sources, tags, and contributions in a named volume — durable across restarts.",
        "Redis is the Celery broker, idle until you turn the worker on.",
        "Qdrant stores embeddings; off by default (flip EMBEDDINGS_ENABLED=true to enable semantic dedup).",
    ])

    # Right column: ASCII-style diagram in code panel
    diagram = (
        "host :50002\n"
        "      │\n"
        "      ▼\n"
        "  ┌────────┐\n"
        "  │ nginx  │\n"
        "  └────┬───┘\n"
        "       │   internal docker network\n"
        "       ▼\n"
        "  ┌────────┐    ┌──────────┐\n"
        "  │FastAPI │ ─► │ Postgres │\n"
        "  │ + CLI  │ ─► │  Redis   │\n"
        "  └────────┘ ─► │  Qdrant  │\n"
        "                └──────────┘"
    )
    add_code_block(s, Inches(7.6), Inches(1.5), Inches(5.1), Inches(4.8), diagram)
    page_number(s, 2, total)


def slide_search_agent(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "2 · LLM-driven search agent")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
             "User types in plain English. Qwen3.5 returns structured filters. The UI applies them as removable chips.",
             size=16, color=TEXT)

    flow = (
        '"openai 2026 blog posts, newest first"\n'
        "         │\n"
        "         ▼   POST /api/items/nl-search\n"
        '         │   (extra_body: enable_thinking=false)\n'
        "         │\n"
        "{\n"
        '  "content_type": "lab_announcement",  ← AI-lab posts ≠ "blog"\n'
        '  "source":       "openai",\n'
        '  "year":         2026,\n'
        '  "sort":         "date_desc"\n'
        "}"
    )
    add_code_block(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(3.2), flow)

    add_bullets(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.7), [
        "Few-shot prompt with taxonomy hints: OpenAI / DeepMind / Meta AI ⇒ lab_announcement, not blog.",
        "Defensive parsing: strips </think> blocks; falls back to plain title-substring search if the LLM emits all-null JSON.",
        "No hardcoded UI dropdowns — the agent owns the search surface; chips are the only knobs.",
    ])
    page_number(s, 3, total)


def slide_ingest(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "3 · Parallel ingest pipeline")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.5),
             "47 sources × ~85 items × 1 LLM call/item is wall-clock bound on the LLM, not CPU. Solution: thread-per-item.",
             size=15, color=TEXT)

    code = (
        "def ingest_source(db, source, workers=settings.ingest_workers):\n"
        "    raw_items = adapter.fetch()              # one HTTP per source\n"
        "    new_ids   = persist_and_dedup(db, raw_items)\n"
        "    db.commit()                              # publish before fan-out\n"
        "    with ThreadPoolExecutor(max_workers=workers) as ex:\n"
        "        list(ex.map(_enrich_one, new_ids))   # one session per worker"
    )
    add_code_block(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.1), code)

    # Knobs panel
    add_panel(s, Inches(0.6), Inches(4.3), Inches(7.5), Inches(2.6))
    add_text(s, Inches(0.85), Inches(4.4), Inches(7.0), Inches(0.4),
             "Knobs", size=15, bold=True, color=DARK)
    add_bullets(s, Inches(0.85), Inches(4.8), Inches(7.0), Inches(2.0), [
        "ingest_workers  =  min(32, cpu // 2)  — per-item LLM enrichment",
        "ingest_source_workers  =  1  — source-level fan-out",
        "DB pool_size  =  max(20, workers + 8)  — one session per worker, no contention",
    ], size=14)

    # Result panel
    add_panel(s, Inches(8.4), Inches(4.3), Inches(4.3), Inches(2.6),
              bg=RGBColor(0xFE, 0xF3, 0xC7), border=AMBER)
    add_text(s, Inches(8.65), Inches(4.4), Inches(3.9), Inches(0.4),
             "Result", size=15, bold=True, color=DARK)
    add_text(s, Inches(8.65), Inches(4.85), Inches(3.9), Inches(2.0),
             "47 sources\n3,889 fetched\n3,516 new\n373 dup\n7 errors\n\n~7 minutes end-to-end\n(32 LLM threads on a 256-core host)",
             size=13, color=TEXT)
    page_number(s, 4, total)


def slide_contribute(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "4 · Contribute & dedup")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.5),
             "Anyone can paste a URL. Qwen reads it; the system classifies, dedups, and the item lands in the feed.",
             size=15, color=TEXT)

    # Pipeline column
    add_text(s, Inches(0.6), Inches(2.0), Inches(6.8), Inches(0.4),
             "Pipeline", size=16, bold=True, color=RED)
    add_bullets(s, Inches(0.6), Inches(2.45), Inches(6.8), Inches(4.5), [
        "Validate URL (scheme + host).",
        "Fetch HTML with project User-Agent + 20s timeout.",
        "Pick the longer of <title> / og:title (≥ 2 words, > 5 chars — kills false dedup matches).",
        "Three-stage dedup: ① canonical-URL match  ② title token_set_ratio ≥ 0.90 (7-day window)  ③ embedding cosine ≥ 0.92 (when enabled).",
        "LLM classifies → topics + content_type + tags.",
        "Insert as Item under built-in 'User contributions' source.",
    ], size=13)

    # Failure UX column
    add_text(s, Inches(7.6), Inches(2.0), Inches(5.1), Inches(0.4),
             "Failure UX", size=16, bold=True, color=RED)
    err = (
        "422 {\n"
        '  "detail": {\n'
        '    "message": "Server returned\n'
        '                HTTP 404.",\n'
        '    "hint": "We need a publicly\n'
        '             accessible HTML page —\n'
        '             not a paywalled URL."\n'
        "  }\n"
        "}"
    )
    add_code_block(s, Inches(7.6), Inches(2.45), Inches(5.1), Inches(3.3), err)
    add_text(s, Inches(7.6), Inches(5.85), Inches(5.1), Inches(1.0),
             "The modal renders message + hint inline so the user knows how to fix the input — no opaque 500s.",
             size=13, italic=True, color=MUTED)
    page_number(s, 5, total)


def slide_tutorial_deploy(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "5 · Tutorial — Deploy in one shot")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.4),
             "Prereqs: Docker 20+ with the compose-v2 plugin.",
             size=15, color=MUTED, italic=True)

    code = (
        "git clone https://github.com/llmsc-security/Hotpot-Tech-Feed.git\n"
        "cd Hotpot-Tech-Feed\n"
        "cp .env.example .env       # set OPENAI_API_KEY, HOST_PORT (default 50002)\n"
        "bash start.sh              # builds images, brings up the stack"
    )
    add_code_block(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.6), code)

    add_text(s, Inches(0.6), Inches(3.8), Inches(12.1), Inches(0.4),
             "start.sh is idempotent: rebuilds, comes up on an internal network, waits for /healthz, prints URLs.",
             size=14, color=TEXT)

    output = (
        "  Hotpot Tech Feed is running.\n"
        "\n"
        "  Open in browser →  http://127.0.0.1:50002\n"
        "  API docs        →  http://127.0.0.1:50002/docs\n"
        "  Health          →  http://127.0.0.1:50002/healthz"
    )
    add_code_block(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(1.5), output)

    add_text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.4),
             "First-time backfill:",
             size=14, bold=True, color=DARK)
    backfill = (
        "docker compose run --rm backend hotpot ingest-now              # default cpu/2 workers\n"
        "docker compose run --rm backend hotpot ingest-now --workers 16 # override"
    )
    add_code_block(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(1.0), backfill)
    page_number(s, 6, total)


def slide_tutorial_migrate(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "6 · Tutorial — Migrate to another PC")

    add_text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
             "Volumes are machine-local. backup.sh / restore.sh turn the live state into a single ~2 MB archive that drops cleanly into any new host.",
             size=14, color=TEXT)

    add_text(s, Inches(0.6), Inches(2.15), Inches(12.1), Inches(0.4),
             "Backup (old PC)",
             size=15, bold=True, color=RED)
    bak = (
        "bash backup.sh\n"
        "# → backups/hotpot-20260430-044721Z.tar.gz\n"
        "#   contains: postgres.dump (pg_dump -Fc), qdrant.tar.gz,\n"
        "#             env.backup, manifest.json"
    )
    add_code_block(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.7), bak)

    add_text(s, Inches(0.6), Inches(4.45), Inches(12.1), Inches(0.4),
             "Restore (new PC)",
             size=15, bold=True, color=RED)
    res = (
        "git clone https://github.com/llmsc-security/Hotpot-Tech-Feed.git\n"
        "cd Hotpot-Tech-Feed\n"
        "bash start.sh                              # empty stack, all defaults\n"
        "scp old-host:.../hotpot-*.tar.gz .\n"
        "bash restore.sh hotpot-20260430-044721Z.tar.gz\n"
        "# stops backend → drops & recreates DB → pg_restore →\n"
        "# restarts → prints /api/stats"
    )
    add_code_block(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.1), res)

    add_text(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.4),
             "No data loss. Postgres dump is logical (cross-version safe); Qdrant snapshot is the raw volume tarball.",
             size=12, color=MUTED, italic=True)
    page_number(s, 7, total)


def slide_closing(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), prs.slide_height - Inches(0.4),
                             prs.slide_width, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.0),
             "Try it", size=54, bold=True, color=DARK, align="center")
    add_text(s, Inches(0.6), Inches(2.9), Inches(12.1), Inches(0.6),
             "http://127.0.0.1:50002",
             size=30, bold=True, color=RED, align="center")

    add_bullets(s, Inches(2.5), Inches(4.0), Inches(8.3), Inches(2.5), [
        'Ask Hotpot  ·  "ML papers from arxiv this year, newest first"',
        "Contribute  ·  paste any blog post URL",
        "Corpus  ·  click the counter → drawer of all sources",
    ], size=18)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
             "feed.ai2wj.com — daily CS digest · powered by Qwen3.5, a free self-hosted LLM",
             size=12, color=MUTED, align="center", italic=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = 8
    slide_title(prs)
    slide_architecture(prs, total)
    slide_search_agent(prs, total)
    slide_ingest(prs, total)
    slide_contribute(prs, total)
    slide_tutorial_deploy(prs, total)
    slide_tutorial_migrate(prs, total)
    slide_closing(prs, total)

    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
