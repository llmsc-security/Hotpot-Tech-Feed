"""Build docs/slides_security.pptx.

Run:
  python3 docs/build_security_slides.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "slides_security.pptx"

DARK = RGBColor(0x0F, 0x17, 0x2A)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xB9, 0x1C, 0x1C)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
BLUE = RGBColor(0x25, 0x63, 0xEB)
PANEL = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CODE_BG = RGBColor(0x0F, 0x17, 0x2A)
CODE_FG = RGBColor(0xE2, 0xE8, 0xF0)


def add_text(slide, left, top, width, height, text, *, size=18, color=INK,
             bold=False, align="left"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}[align]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, text, subtitle=None):
    add_text(slide, Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.65),
             text, size=31, bold=True, color=DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.08),
                                 Inches(1.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.2), Inches(12.1), Inches(0.4),
                 subtitle, size=14, color=MUTED)


def add_panel(slide, left, top, width, height, *, bg=PANEL, border=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(0.8)
    shape.shadow.inherit = False
    return shape


def add_badge(slide, left, top, text, *, color=BLUE, width=1.55):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   Inches(width), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Pt(5)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = 2
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_bullets(slide, left, top, width, height, items, *, size=15, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = "- " + item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_code(slide, left, top, width, height, text, *, size=13):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = CODE_BG
    panel.line.fill.background()
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(12)
    tf.margin_top = tf.margin_bottom = Pt(10)
    for i, line in enumerate(text.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def page(slide, n, total):
    add_text(slide, Inches(12.1), Inches(7.02), Inches(0.7), Inches(0.25),
             f"{n}/{total}", size=10, color=MUTED, align="right")


def title_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(0.38), prs.slide_height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RED
    stripe.line.fill.background()
    add_text(s, Inches(0.75), Inches(1.65), Inches(11.8), Inches(0.8),
             "Hotpot /security", size=42, bold=True, color=DARK)
    add_text(s, Inches(0.75), Inches(2.55), Inches(11.8), Inches(0.55),
             "Evidence-first scoring, ranking, and anti-marketing algorithm",
             size=22, color=RED)
    add_text(s, Inches(0.75), Inches(3.55), Inches(11.2), Inches(1.0),
             "Goal: surface real CVEs, real exploitation, concrete cases, and actionable reports. "
             "The general feed keeps Item.score; /security owns final_security_score.",
             size=18, color=INK)
    add_text(s, Inches(0.75), Inches(5.65), Inches(11.2), Inches(0.45),
             "Current v1: deterministic local article/source signals, persisted in security_item_scores.",
             size=15, color=MUTED)
    page(s, 1, total)


def evidence_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "1. Evidence collection", "Convert each item into observable facts before any ranking.")

    cols = [
        ("Identity", ["CVE IDs", "GHSA IDs", "CWE IDs", "affected products", "affected versions"]),
        ("Exploitation", ["CISA KEV signal", "confirmed in the wild", "vendor confirmed", "credible incident report", "public PoC"]),
        ("Actionability", ["patch available", "patched version", "mitigation", "detection logic", "IoCs / YARA / Sigma"]),
        ("Quality / noise", ["source links", "technical detail", "generic phrases", "demo / webinar CTA", "product pitch"]),
    ]
    x = 0.55
    for title, bullets in cols:
        add_panel(s, Inches(x), Inches(1.75), Inches(3.0), Inches(4.45))
        add_text(s, Inches(x + 0.18), Inches(1.95), Inches(2.65), Inches(0.35),
                 title, size=16, bold=True, color=DARK)
        add_bullets(s, Inches(x + 0.18), Inches(2.45), Inches(2.65), Inches(3.45),
                    bullets, size=13)
        x += 3.15

    add_text(s, Inches(0.65), Inches(6.45), Inches(12.0), Inches(0.45),
             "Important rule: the LLM may extract facts later, but ranking is deterministic code.",
             size=14, color=MUTED)
    page(s, 2, total)


def pipeline_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "2. Filter first, rank second", "Bad security-looking content should not enter the candidate set.")

    code = (
        "raw crawled item\n"
        "  -> feature extraction\n"
        "  -> security relevance filter     threshold >= 0.40\n"
        "  -> soft-article rejection        marketing + weak evidence\n"
        "  -> evidence score                concrete, source-backed facts\n"
        "  -> exploitation score            KEV / active exploitation / PoC ladder\n"
        "  -> quality, impact, actionability\n"
        "  -> grouping by CVE / GHSA / incident\n"
        "  -> final_security_score + security_hot_score\n"
        "  -> /api/security/hot and /api/security/items"
    )
    add_code(s, Inches(0.7), Inches(1.75), Inches(6.5), Inches(4.65), code, size=14)

    add_panel(s, Inches(7.55), Inches(1.75), Inches(4.95), Inches(4.65),
              bg=RGBColor(0xFE, 0xF2, 0xF2), border=RGBColor(0xFE, 0xCA, 0xCA))
    add_text(s, Inches(7.85), Inches(1.98), Inches(4.35), Inches(0.4),
             "Reject examples", size=17, bold=True, color=RED)
    add_bullets(s, Inches(7.85), Inches(2.48), Inches(4.25), Inches(2.1), [
        "Top 10 cybersecurity trends",
        "Why every CISO needs our platform",
        "Download our latest report",
        "Best practices with no CVE, patch, IoC, victim, or actor",
    ], size=13)
    add_text(s, Inches(7.85), Inches(4.85), Inches(4.25), Inches(0.85),
             "Promotional language alone is not enough to reject. Promotional language plus weak evidence is rejected.",
             size=13, color=INK)
    page(s, 3, total)


def formula_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "3. Final security score", "Separate from the general feed Item.score.")

    formula = (
        "final_security_score =\n"
        "    0.30 * evidence_score\n"
        "  + 0.24 * exploitation_score\n"
        "  + 0.14 * content_quality_score\n"
        "  + 0.10 * impact_score\n"
        "  + 0.08 * actionability_score\n"
        "  + 0.06 * source_authority_score\n"
        "  + 0.05 * corroboration_score\n"
        "  + 0.03 * freshness_score\n"
        "  - 0.22 * soft_article_score\n\n"
        "Clamp to [0, 1]"
    )
    add_code(s, Inches(0.7), Inches(1.72), Inches(6.25), Inches(4.2), formula, size=14)

    add_text(s, Inches(7.35), Inches(1.8), Inches(4.9), Inches(0.35),
             "Priority order", size=17, bold=True, color=DARK)
    add_bullets(s, Inches(7.35), Inches(2.28), Inches(4.9), Inches(2.3), [
        "Evidence and exploitation dominate.",
        "Quality, impact, and actionability decide usefulness.",
        "Authority, corroboration, and freshness help tie-break.",
        "Soft-article score can push weak promotional items out.",
    ], size=14)

    add_panel(s, Inches(7.35), Inches(4.85), Inches(4.9), Inches(1.1),
              bg=RGBColor(0xEC, 0xFE, 0xFF), border=RGBColor(0xA5, 0xF3, 0xFC))
    add_text(s, Inches(7.62), Inches(5.05), Inches(4.35), Inches(0.6),
             "Freshness and clicks cannot rescue weak evidence.",
             size=16, bold=True, color=DARK)
    page(s, 4, total)


def exploitation_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "4. Exploitation and evidence ladders", "PoC is useful, but it is not proof of real-world exploitation.")

    add_panel(s, Inches(0.65), Inches(1.72), Inches(5.75), Inches(4.8))
    add_text(s, Inches(0.9), Inches(1.95), Inches(5.25), Inches(0.35),
             "Exploitation score", size=17, bold=True, color=DARK)
    ladder = [
        ("1.00", "CISA KEV"),
        ("0.90", "confirmed in the wild"),
        ("0.85", "vendor-confirmed exploitation"),
        ("0.70", "credible report claims exploitation"),
        ("0.45", "public PoC available"),
        ("0.20", "theoretical vulnerability only"),
        ("0.00", "unknown"),
    ]
    y = 2.48
    for score, label in ladder:
        add_badge(s, Inches(0.9), Inches(y), score, color=RED if score in {"1.00", "0.90", "0.85"} else BLUE, width=0.78)
        add_text(s, Inches(1.85), Inches(y + 0.02), Inches(4.25), Inches(0.3),
                 label, size=13, color=INK)
        y += 0.48

    add_panel(s, Inches(6.75), Inches(1.72), Inches(5.75), Inches(4.8))
    add_text(s, Inches(7.0), Inches(1.95), Inches(5.25), Inches(0.35),
             "Evidence score rewards", size=17, bold=True, color=DARK)
    add_bullets(s, Inches(7.0), Inches(2.45), Inches(5.15), Inches(3.55), [
        "validated CVE / GHSA identity",
        "KEV, vendor advisory, primary research, credible media",
        "patch, mitigation, affected versions, IoCs",
        "victim, threat actor, timeline, campaign detail",
        "source links and repeat exposure",
    ], size=14)
    page(s, 5, total)


def acceptance_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "5. Acceptance rule and anti-marketing filter", "The /security page is a high-precision projection.")

    accept = (
        "accept_item =\n"
        "    security_relevance_score >= 0.40\n"
        "    AND (\n"
        "        evidence_score >= 0.35\n"
        "        OR cisa_kev_match = true\n"
        "        OR confirmed_exploitation = true\n"
        "    )\n"
        "    AND soft_article_reject = false\n"
        "    AND final_security_score >= 0.30"
    )
    add_code(s, Inches(0.75), Inches(1.75), Inches(6.3), Inches(3.45), accept, size=13)

    add_panel(s, Inches(7.4), Inches(1.75), Inches(4.85), Inches(3.45),
              bg=RGBColor(0xFF, 0xFB, 0xEB), border=RGBColor(0xFD, 0xBA, 0x74))
    add_text(s, Inches(7.7), Inches(1.98), Inches(4.3), Inches(0.35),
             "Soft-article reject", size=17, bold=True, color=RED)
    add_bullets(s, Inches(7.7), Inches(2.45), Inches(4.25), Inches(2.2), [
        "score >= 0.75",
        "score >= 0.55 and evidence < 0.45",
        "no CVE/GHSA/victim/actor/patch/IoC and score >= 0.45",
    ], size=13)

    add_text(s, Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.6),
             "This prevents a broad corpus from leaking marketing posts into the security feed.",
             size=17, bold=True, color=DARK)
    page(s, 6, total)


def ranking_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "6. Group-level ranking", "The page displays story groups, not raw duplicate items.")

    add_code(s, Inches(0.7), Inches(1.68), Inches(5.9), Inches(2.65), (
        "group_final_security_score =\n"
        "    max(item.final_security_score)\n"
        "  + min(0.06, 0.015 * additional_authoritative_sources)\n\n"
        "group_security_hot_score =\n"
        "    max(item.security_hot_score)\n"
        "  + min(0.05, 0.0125 * additional_recent_sources)"
    ), size=12)

    add_code(s, Inches(6.95), Inches(1.68), Inches(5.6), Inches(2.65), (
        "Top 10 hot sort:\n"
        "  security_hot_score desc\n"
        "  final_security_score desc\n"
        "  exploitation_score desc\n"
        "  evidence_score desc\n"
        "  event_time desc\n\n"
        "Pager sort:\n"
        "  final_security_score desc\n"
        "  event_time desc\n"
        "  group_key asc"
    ), size=12)

    add_panel(s, Inches(0.7), Inches(4.72), Inches(11.85), Inches(1.25),
              bg=RGBColor(0xEC, 0xFD, 0xF5), border=RGBColor(0xA7, 0xF3, 0xD0))
    add_text(s, Inches(1.0), Inches(4.93), Inches(11.25), Inches(0.72),
             "Representative card: highest final_security_score; tie-break by source authority, then newest event_time.",
             size=17, bold=True, color=DARK)
    page(s, 7, total)


def api_slide(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "7. API, persistence, and current result", "Backend owns scores so pagination and top 10 are stable.")

    add_panel(s, Inches(0.65), Inches(1.7), Inches(5.8), Inches(4.5))
    add_text(s, Inches(0.9), Inches(1.95), Inches(5.3), Inches(0.35),
             "Persisted projection", size=17, bold=True, color=DARK)
    add_bullets(s, Inches(0.9), Inches(2.45), Inches(5.25), Inches(2.75), [
        "security_item_scores table",
        "final_security_score and security_hot_score",
        "badges, why_ranked, source_chain",
        "event_time, section, group_key",
        "score_version = security-v1",
    ], size=14)
    add_text(s, Inches(0.9), Inches(5.35), Inches(5.25), Inches(0.45),
             "Backfill: 7,007 processed, 109 accepted, 6,898 rejected, 0 errors.",
             size=13, color=MUTED)

    add_panel(s, Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.5))
    add_text(s, Inches(7.1), Inches(1.95), Inches(5.3), Inches(0.35),
             "Read contracts", size=17, bold=True, color=DARK)
    add_code(s, Inches(7.1), Inches(2.45), Inches(5.25), Inches(2.15), (
        "GET /api/security/hot?limit=10\n\n"
        "GET /api/security/items\n"
        "  ?limit=25\n"
        "  &offset=0\n"
        "  &section=all\n"
        "  &sort=score_desc"
    ), size=12)
    add_text(s, Inches(7.1), Inches(4.9), Inches(5.25), Inches(0.55),
             "UI: two rows of five hot cards on desktop, pager after filtering and grouping.",
             size=13, color=MUTED)
    page(s, 8, total)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 8
    title_slide(prs, total)
    evidence_slide(prs, total)
    pipeline_slide(prs, total)
    formula_slide(prs, total)
    exploitation_slide(prs, total)
    acceptance_slide(prs, total)
    ranking_slide(prs, total)
    api_slide(prs, total)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
